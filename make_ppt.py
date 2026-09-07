import ollama
from pptx import Presentation
from pptx.util import Pt

# 1. 设定一个需要深度思考的主题
topic = "2026年半导体硬件供应链投资机遇与基本面分析"

# 2. 用“专业身份 + 宽松格式”去引导大模型
prompt = f"""
你是一位顶级的券商投行分析师。
我需要你为主题《{topic}》构思一份有深度、有洞见的 PPT 大纲，总共 3 页。
请自由发挥你的专业知识，深入思考行业基本面、价格行为趋势和财务估值，写出极具商业价值的硬核文案。

【输出规则】（为了让我的程序能读取，请务必遵守）：
1. 每一页的幻灯片内容之间，请用连续的三个等号 "===" 隔开。
2. 每一页里，第一行必须是标题，第二行开始是你深度思考的正文要点。
3. 不要输出任何寒暄或解释性的废话，直接给我最终的内容。

示例格式：
核心科技公司的资产重估
- 随着设备折旧周期的结束，部分半导体企业的利润率将迎来拐点。
- 结合近期的价格行为分析，供应链上游的零部件厂商具备更高的安全边际。
===
下一页的标题
- 洞见1...
- 洞见2...
"""

print("正在调用本地大模型进行深度思考与文案创作...")
response = ollama.chat(
    model='llama3.1', # 确保这是你本地的模型名称
    messages=[{'role': 'user', 'content': prompt}]
)

ai_output = response['message']['content']
print("文案创作完毕，正在自动排版 PPT...\n")

# 打印出大模型的原始回答，方便你在终端里查看它的思考质量
print("【大模型生成的文案如下】：")
print(ai_output)
print("\n------------------------\n")

# 3. 创建并组装 PPT
prs = Presentation()

# 使用 === 来切分每一页的内容（非常稳定的切分方式）
slides_content = ai_output.split('===')

for slide_text in slides_content:
    slide_text = slide_text.strip()
    if not slide_text:
        continue # 如果是空的就跳过
    
    # 按行把这一页的内容切开
    lines = [line.strip() for line in slide_text.split('\n') if line.strip()]
    
    if len(lines) >= 1:
        # 第一行当标题
        title = lines[0].replace("#", "").replace("*", "").strip() 
        # 剩下的所有行当正文
        body_points = lines[1:] 
        
        # 添加幻灯片
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 写入标题
        slide.shapes.title.text = title
        
        # 写入正文（支持无限多行的大模型自由发挥）
        if body_points:
            tf = slide.placeholders[1].text_frame
            tf.text = body_points[0].replace("-", "").replace("*", "").strip()
            for point in body_points[1:]:
                p = tf.add_paragraph()
                p.text = point.replace("-", "").replace("*", "").strip()
                p.font.size = Pt(18)

# 4. 保存文件
output_file = "C:/Users/Lyang/Desktop/深度洞见_汇报演示文稿.pptx"
prs.save(output_file)

print(f"🎉 成功！更聪明的 PPT 已制作完成并保存至桌面: {output_file}")